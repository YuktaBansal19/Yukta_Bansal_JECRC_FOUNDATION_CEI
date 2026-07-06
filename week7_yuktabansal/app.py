import streamlit as st
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter  
import os
import io
from pptx import Presentation
from PIL import Image
import pytesseract
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_core.prompts import PromptTemplate 
import shutil 
from dotenv import load_dotenv


load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
try:
    api_key = api_key or st.secrets.get("GOOGLE_API_KEY")
except:
    pass
if not api_key:
    st.error("GOOGLE_API_KEY not found!")
    st.stop()
    
genai.configure(api_key=api_key)


if shutil.which("tesseract"):
    pytesseract.pytesseract.tesseract_cmd = shutil.which("tesseract")

def get_text(file):
    text = ""
    file_extension = os.path.splitext(file.name)[1].lower()
    try:
        if file_extension == ".pdf":
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
 
        elif file_extension == ".pptx":
            prs = Presentation(io.BytesIO(file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
 
        elif file_extension in (".txt", ".md"):
            text = file.read().decode("utf-8")
 
        elif file_extension in (".jpg", ".jpeg", ".png"):
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
 
        else:
            st.warning(f"Unsupported file format: {file_extension}")
 
    except Exception as e:
        st.error(f"Error processing {file.name}: {str(e)}")
 
    return text
 
 
def get_files_text(files):
    """Process multiple files and combine extracted text."""
    text = ""
    for file in files:
        extracted = get_text(file)
        if extracted.strip():
            text += extracted + "\n\n"
        else:
            st.warning(f"No text extracted from: {file.name}")
    return text
 
 
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000, chunk_overlap=200
    )
    return text_splitter.split_text(text)
 
 
def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
        return True
    except Exception as e:
        st.error(f"Error creating vector store: {str(e)}")
        return False
 
 
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context.
    If the answer is not in the context, say "answer is not available in the context".

    Context:\n{context}\n
    Question:\n{question}\n
    Answer:
    """
    try:
        model = ChatGoogleGenerativeAI(model="models/gemini-2.5-flash", temperature=0.3)
        prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
        chain = prompt | model | StrOutputParser()
        return chain
    except Exception as e:
        st.error(f"Error creating conversational chain: {str(e)}")
        return None


def user_input(user_question):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")

        save_path = os.path.join(os.path.dirname(__file__), "faiss_index")
        if not os.path.exists(save_path):
            st.error("Please upload and process your files first!")
            return

        new_db = FAISS.load_local(save_path, embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)

        # combine docs into a single context string
        context = "\n\n".join([doc.page_content for doc in docs])

        chain = get_conversational_chain()
        if chain is None:
            return

        response = chain.invoke({"context": context, "question": user_question})
        st.write("**Reply:** ", response)

    except Exception as e:
        st.error(f"Error processing question: {str(e)}")
 
 
def main():
    st.set_page_config(page_title="MultiDocumentChatBot", layout="wide")
    st.header("Ask-My-Docs")
 
    user_question = st.text_input("Ask a Question from the Files")
    if user_question:
        user_input(user_question)
 
    with st.sidebar:
        st.title("Menu:")
        files = st.file_uploader(
            "Upload your Files and Click on Submit & Process",
            accept_multiple_files=True,
            type=["pdf", "pptx", "txt", "jpg", "jpeg", "png", "md"],
        )
 
        if st.button("Submit & Process"):
            if not files:
                st.error("Please upload at least one file!")
            else:
                with st.spinner("Processing..."):
                    try:
                        raw_text = get_files_text(files)
                        if not raw_text.strip():
                            st.error("No text could be extracted from the uploaded files!")
                            return
                        text_chunks = get_text_chunks(raw_text)
                        if get_vector_store(text_chunks):
                            st.success("Files processed! You can now ask questions.")
                        else:
                            st.error("Failed to process files!")
                    except Exception as e:
                        st.error(f"Error processing files: {str(e)}")
 
if __name__ == "__main__":
    main()